import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR

class WordToPPTConverter:
    def __init__(self, root):
        self.root = root
        self.root.title("Word转PPT转换器")
        self.setup_ui()
        
    def setup_ui(self):
        # 输入文件选择
        tk.Label(self.root, text="输入Word文件:").grid(row=0, column=0, sticky="w")
        self.input_entry = tk.Entry(self.root, width=40)
        self.input_entry.grid(row=0, column=1)
        tk.Button(self.root, text="浏览...", command=self.select_input_file).grid(row=0, column=2)
        
        # 输出文件选择
        tk.Label(self.root, text="输出PPT文件:").grid(row=1, column=0, sticky="w")
        self.output_entry = tk.Entry(self.root, width=40)
        self.output_entry.grid(row=1, column=1)
        tk.Button(self.root, text="浏览...", command=self.select_output_file).grid(row=1, column=2)
        
        # 幻灯片比例（只保留下拉式選單）
        tk.Label(self.root, text="幻灯片比例:").grid(row=2, column=0, sticky="w")
        ratios = ["4:3", "16:9", "9:18"]
        self.aspect_ratio = tk.StringVar(value=ratios[1])
        tk.OptionMenu(self.root, self.aspect_ratio, *ratios).grid(row=2, column=1)
        
        # 字体设置（下拉式選單，显示系統所有中文字型）
        tk.Label(self.root, text="字体:").grid(row=3, column=0, sticky="w")
        
        # 获取系统中所有可用字体
        import tkinter.font as tkFont
        all_fonts = list(tkFont.families())
        
        # 筛选出包含中文关键字的字体
        chinese_keywords = ["宋", "黑", "楷", "仿", "雅", "明", "正", "宅", "標", "細", "體"]
        chinese_fonts = [font for font in all_fonts 
                        if any(keyword in font for keyword in chinese_keywords)]
        
        # 如果没有找到中文字体，使用默认列表
        if not chinese_fonts:
            chinese_fonts = ["微軟正黑", "標楷體", "新細明體", "黑體", "微软雅黑", "宋体", "仿宋", "楷体"]
            
        self.font_family = tk.StringVar(value=chinese_fonts[0] if chinese_fonts else "")
        tk.OptionMenu(self.root, self.font_family, *chinese_fonts).grid(row=3, column=1)
        
        tk.Label(self.root, text="字体大小:").grid(row=4, column=0, sticky="w")
        font_sizes = [16, 18, 20, 22, 24, 28, 32, 36, 40, 48, 56, 64, 72]
        self.font_size = tk.IntVar(value=24)
        tk.OptionMenu(self.root, self.font_size, *font_sizes).grid(row=4, column=1)
        
        tk.Label(self.root, text="字体颜色:").grid(row=5, column=0, sticky="w")
        self.font_color = tk.StringVar(value="#000000")
        tk.Entry(self.root, textvariable=self.font_color, width=10).grid(row=5, column=1, sticky="w")
        tk.Button(self.root, text="選擇字體顏色", command=self.choose_font_color).grid(row=5, column=2, sticky="w")
        self.font_color_block = tk.Label(self.root, width=2, height=1, bg=self.font_color.get())
        self.font_color_block.grid(row=5, column=3, padx=5)
        self.font_color.trace_add("write", lambda *args: self.update_color_blocks())

        tk.Label(self.root, text="背景颜色:").grid(row=6, column=0, sticky="w")
        self.bg_color = tk.StringVar(value="#FFFFFF")
        tk.Entry(self.root, textvariable=self.bg_color, width=10).grid(row=6, column=1, sticky="w")
        tk.Button(self.root, text="選擇背景顏色", command=self.choose_bg_color).grid(row=6, column=2, sticky="w")
        self.bg_color_block = tk.Label(self.root, width=2, height=1, bg=self.bg_color.get())
        self.bg_color_block.grid(row=6, column=3, padx=5)
        self.bg_color.trace_add("write", lambda *args: self.update_color_blocks())
        
        # 每行字数设置
        tk.Label(self.root, text="每行字数:").grid(row=7, column=0, sticky="w")
        chars_per_line = list(range(15, 31))  # 15到30
        self.chars_per_line = tk.IntVar(value=20)
        tk.OptionMenu(self.root, self.chars_per_line, *chars_per_line).grid(row=7, column=1)
        
        # 每页行数设置
        tk.Label(self.root, text="每页行数:").grid(row=8, column=0, sticky="w")
        lines_per_page = [6, 8, 10, 12]
        self.lines_per_page = tk.IntVar(value=8)
        tk.OptionMenu(self.root, self.lines_per_page, *lines_per_page).grid(row=8, column=1)
        
        # 转换按钮（调整到第9行）
        tk.Button(self.root, text="开始转换", command=self.convert, bg="green", fg="white").grid(row=9, column=1, pady=10)
    
    def select_input_file(self):
        filepath = filedialog.askopenfilename(filetypes=[("Word文件", "*.docx")])
        if filepath:
            self.input_entry.delete(0, tk.END)
            self.input_entry.insert(0, filepath)
    
    def select_output_file(self):
        filepath = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PPT文件", "*.pptx")])
        if filepath:
            self.output_entry.delete(0, tk.END)
            self.output_entry.insert(0, filepath)
    
    def hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip('#')
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def convert(self):
        try:
            # 讀取Word文檔
            doc = Document(self.input_entry.get())
            
            # 自動命名功能
            import os
            from datetime import datetime
            input_path = self.input_entry.get()
            output_path = self.output_entry.get()
            if not output_path.strip():
                now = datetime.now()
                mmss = now.strftime("%M_%S")
                word_name = os.path.splitext(os.path.basename(input_path))[0]
                word_dir = os.path.dirname(input_path)
                output_path = os.path.join(word_dir, f"{mmss}_{word_name}.pptx")
                self.output_entry.delete(0, tk.END)
                self.output_entry.insert(0, output_path)
    
            # 創建PPT
            prs = Presentation()
    
            # 設定幻燈片比例
            ratio = self.aspect_ratio.get()
            if ratio == "4:3":
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
            elif ratio == "16:9":
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(5.625)
            elif ratio == "9:18":
                prs.slide_width = Inches(5)
                prs.slide_height = Inches(10)
    
            # 取得顏色設定
            font_rgb = self.hex_to_rgb(self.font_color.get())
            bg_rgb = self.hex_to_rgb(self.bg_color.get())
    
            # 合併所有段落並處理換行
            all_text = []
            for paragraph in doc.paragraphs:
                # 每20字插入換行符
                text = paragraph.text
                processed_text = ""
                for i in range(0, len(text), 20):
                    processed_text += text[i:i+20] + "\n"
                all_text.append(processed_text.strip())
            
            # 合併所有處理後的文本
            merged_text = "\n".join(all_text)
            
            # 处理每行字数
            processed_lines = []
            for line in merged_text.splitlines():
                # 按设置的字数分割行
                for i in range(0, len(line), self.chars_per_line.get()):
                    processed_lines.append(line[i:i+self.chars_per_line.get()])
            
            # 按设置的行数分页
            for i in range(0, len(processed_lines), self.lines_per_page.get()):
                page_lines = processed_lines[i:i+self.lines_per_page.get()]
                chunk = "\n".join(page_lines)
                
                # 創建PPT頁面
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide_bg = slide.background
                fill = slide_bg.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*bg_rgb)
                
                left = (prs.slide_width - prs.slide_width * 0.85) / 2
                top = (prs.slide_height - prs.slide_height * 0.85) / 2
                txBox = slide.shapes.add_textbox(left, top, prs.slide_width * 0.85, prs.slide_height * 0.85)
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                p = tf.add_paragraph()
                p.text = chunk
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = self.font_family.get()
                    run.font.size = Pt(self.font_size.get())
                    run.font.color.rgb = RGBColor(*font_rgb)
                fill = txBox.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*bg_rgb)
    
            # 保存PPT
            prs.save(self.output_entry.get())
            messagebox.showinfo("成功", "转换完成！")
        except Exception as e:
            messagebox.showerror("错误", f"转换失败: {str(e)}")

    def update_color_blocks(self):
        try:
            self.font_color_block.config(bg=self.font_color.get())
        except Exception:
            pass
        try:
            self.bg_color_block.config(bg=self.bg_color.get())
        except Exception:
            pass

    def choose_font_color(self):
        from tkinter import colorchooser
        color = colorchooser.askcolor(title="選擇字體顏色")
        if color[1]:
            self.font_color.set(color[1])
        self.update_color_blocks()

    def choose_bg_color(self):
        from tkinter import colorchooser
        color = colorchooser.askcolor(title="選擇背景顏色")
        if color[1]:
            self.bg_color.set(color[1])
        self.update_color_blocks()

if __name__ == "__main__":
    root = tk.Tk()
    app = WordToPPTConverter(root)
    root.mainloop()