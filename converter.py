import os
import time
import threading
import subprocess
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
import tkinter as tk
from tkinter import messagebox

class Converter:
    def __init__(self, app):
        self.app = app

    def start_conversion(self):
        if not self.app.input_path or not self.app.ppt_path:
            messagebox.showerror("錯誤", "請選擇輸入和輸出文件！")
            return

        self.app.convert_btn.config(state=tk.DISABLED)
        self.app.converting = True
        self.app.success_flag = {'ok': True}

        if self.app.input_path.lower().endswith(".pdf"):
            self.app.input_type = "pdf"
        elif self.app.input_path.lower().endswith(".docx"):
            self.app.input_type = "docx"
        elif self.app.input_path.lower().endswith(".txt"):
            self.app.input_type = "txt"
        else:
            self.app.success_flag['ok'] = False
            self.app.root.after(0, lambda: self._show_error("不支援的檔案類型"))
            return

        threading.Thread(target=self._run_conversion).start()
        threading.Thread(target=self._animate_loading).start()

    def _animate_loading(self):
        dots = ""
        while self.app.converting:
            dots = dots + "." if len(dots) < 3 else ""
            self.app.loading_label.config(text=f"轉檔中{dots}")
            self.app.root.update_idletasks()
            time.sleep(0.5)

    def _run_conversion(self):
        pdf_path = None
        doc = None
        try:
            # 根據模板創建簡報
            if self.app.template_path and os.path.exists(self.app.template_path):
                prs = Presentation(self.app.template_path)
            else:
                prs = Presentation()
            
            # 移除模板原有內容（根據需要保留部分內容）
            while len(prs.slides) > 0:
                xml_slides = prs.slides._sldIdLst  
                xml_slides.remove(xml_slides[0])
            
            ratio = self.app.aspect_ratio.get()
            if ratio == "16:9":
                prs.slide_width = Inches(16)
                prs.slide_height = Inches(9)
            elif ratio == "4:3":
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
            elif ratio == "10:16":
                prs.slide_width = Inches(9)
                prs.slide_height = Inches(16)

            # 處理文字文件
            if self.app.input_path.endswith(".txt"):
                self._convert_txt_to_ppt(prs)
            elif self.app.input_type == 'docx':
                pdf_path = self._convert_docx_to_pdf()
                if not pdf_path:
                    return
                self._convert_pdf_to_ppt(prs, pdf_path)
            else:
                self._convert_pdf_to_ppt(prs, self.app.input_path)

            prs.save(self.app.ppt_path)

        except Exception as e:
            self.app.success_flag['ok'] = False
            self.app.root.after(0, lambda err=e: self._show_error(str(err)))
        finally:
            self.app.converting = False
            if doc:
                doc.close()
            if self.app.input_type == 'docx' and pdf_path and os.path.exists(pdf_path):
                try:
                    os.remove(pdf_path)
                except Exception as e:
                    print(f"刪除臨時檔案失敗: {e}")
            self.app.root.after(0, self._finish_success)

    def _convert_txt_to_ppt(self, prs):
        with open(self.app.input_path, "r", encoding="utf-8") as f:
            text = f.read()
        
        # 根據兩個以上空白行分割文字內容
        page_contents = []
        current_page = []
        empty_line_count = 0
        
        for line in text.splitlines():
            if not line.strip():  # 空行
                empty_line_count += 1
            else:
                if empty_line_count >= 2:  # 兩個以上空白行表示換頁
                    if current_page:  # 確保不添加空頁面
                        page_contents.append("\n".join(current_page))
                        current_page = []
                empty_line_count = 0
                current_page.append(line)
        
        # 添加最後一頁
        if current_page:
            page_contents.append("\n".join(current_page))
        
        # 如果沒有內容，創建一個空頁面
        if not page_contents:
            page_contents = [""]
        
        # 為每一頁創建幻燈片
        for page_content in page_contents:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            if self.app.page_bg_color:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*self.app.page_bg_color)
            
            textbox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
            tf = textbox.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            textbox.left = Inches(0)
            textbox.top = Inches(0)
            textbox.width = prs.slide_width
            textbox.height = prs.slide_height
            
            lines = page_content.splitlines()
            
            for line in lines:
                indent = 0
                for char in line:
                    if char == ' ':
                        indent += 1
                    else:
                        break
                level = min(max(indent // 4, 0), 8)
                clean_line = line.lstrip()
                
                if not clean_line:
                    continue
                
                p = tf.add_paragraph()
                p.text = clean_line
                p.level = level
                
                align_value = self.app.text_align.get()
                if align_value == "LEFT":
                    p.alignment = PP_ALIGN.LEFT
                elif align_value == "CENTER":
                    p.alignment = PP_ALIGN.CENTER
                elif align_value == "RIGHT":
                    p.alignment = PP_ALIGN.RIGHT
                
                run = p.runs[0]
                run.font.size = Pt(self.app.font_size)
                run.font.color.rgb = RGBColor(*self.app.font_color)
                run.font.name = self.app.font_name
            
            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

    def _convert_pdf_to_ppt(self, prs, pdf_path):
        doc = fitz.open(pdf_path)
        for page in doc:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            if self.app.page_bg_color:
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(*self.app.page_bg_color)

            textbox = slide.shapes.add_textbox(0, 0, prs.slide_width, prs.slide_height)
            tf = textbox.text_frame
            tf.clear()
            # Set text frame to center align both horizontally and vertically
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # Vertical center
            textbox.left = Inches(0)
            textbox.top = Inches(0)
            textbox.width = prs.slide_width
            textbox.height = prs.slide_height

            text = page.get_text("text", flags=0)
            lines = text.splitlines()

            for line in lines:
                indent = 0
                for char in line:
                    if char == ' ':
                        indent += 1
                    else:
                        break
                level = min(max(indent // 4, 0), 8)  # 限制層級在 0~8 之間 
                clean_line = line.lstrip()
                
                if not clean_line:  # 跳過空行
                    continue
                
                p = tf.add_paragraph()
                p.text = clean_line
                p.level = level
                # 根據選擇的對齊方式設定文字對齊
                align_value = self.app.text_align.get()
                if align_value == "LEFT":
                    p.alignment = PP_ALIGN.LEFT
                elif align_value == "CENTER":
                    p.alignment = PP_ALIGN.CENTER
                elif align_value == "RIGHT":
                    p.alignment = PP_ALIGN.RIGHT

                run = p.runs[0]
                run.font.size = Pt(self.app.font_size)
                run.font.color.rgb = RGBColor(*self.app.font_color)
                run.font.name = self.app.font_name

            tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        doc.close()

    def _convert_docx_to_pdf(self):
        try:
            pdf_path = os.path.splitext(self.app.input_path)[0] + ".pdf"
            result = subprocess.run(
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", os.path.dirname(self.app.input_path), self.app.input_path],
                capture_output=True,
                text=True
            )
            if result.returncode != 0:
                raise Exception(f"LibreOffice 錯誤：\n{result.stderr}")
            return pdf_path
        except FileNotFoundError:
            self.app.root.after(0, lambda: self._show_error("請先安裝 LibreOffice"))
            return None
        except Exception as e:
            self.app.root.after(0, lambda: self._show_error(f"DOCX 轉 PDF 失敗: {str(e)}"))
            return None

    def _finish_success(self):
        self.app.convert_btn.config(state=tk.NORMAL)
        if self.app.success_flag['ok']:
            messagebox.showinfo("成功", "文件轉換完成！")
        self.app.loading_label.config(text="")

    def _show_error(self, msg: str):
        self.app.converting = False
        self.app.loading_label.config(text="轉換失敗 ✘")
        messagebox.showerror("錯誤", f"轉換失敗：{msg}")
        self.app.convert_btn.config(state=tk.NORMAL)